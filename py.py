#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
将Markdown演示文稿转换为PPTX格式
使用 python-pptx 库生成PowerPoint文件
"""

import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def parse_markdown(content):
    """解析Markdown内容，提取幻灯片"""
    slides_data = []
    # 使用 --- 分割幻灯片
    slides = content.split('\n---\n')
    
    for slide in slides:
        slide = slide.strip()
        if not slide or slide.startswith('marp:') or slide.startswith('theme:') or slide.startswith('_class:') or 'style:' in slide:
            continue
        slides_data.append(slide)
    
    return slides_data

def extract_title_and_content(slide_text):
    """从幻灯片文本中提取标题和内容"""
    lines = slide_text.split('\n')
    title = ""
    content = []
    
    for i, line in enumerate(lines):
        line = line.strip()
        if not title and line.startswith('#'):
            # 提取标题
            title = re.sub(r'^#+\s*', '', line)
        elif line and not line.startswith('#'):
            content.append(line)
    
    return title, content

def add_slide_with_title_and_content(prs, title, content_lines):
    """添加包含标题和内容的幻灯片"""
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    
    # 添加标题
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(44)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)
    
    # 添加内容
    body_shape = slide.placeholders[1]
    text_frame = body_shape.text_frame
    text_frame.clear()
    
    current_level = 0
    for line in content_lines:
        if not line.strip():
            continue
        
        # 计算缩进级别
        level = (len(line) - len(line.lstrip())) // 2
        
        # 处理不同的标记类型
        clean_line = line.strip()
        
        # 移除Markdown标记
        clean_line = re.sub(r'^\*\s*', '', clean_line)  # 移除 * 
        clean_line = re.sub(r'^-\s*', '', clean_line)   # 移除 -
        clean_line = re.sub(r'^\d+\.\s*', '', clean_line)  # 移除 1. 2. 等
        clean_line = re.sub(r'^\*\*', '', clean_line)   # 移除 **
        clean_line = re.sub(r'\*\*$', '', clean_line)
        clean_line = re.sub(r'`([^`]+)`', r'\1', clean_line)  # 移除代码标记
        
        if clean_line:
            p = text_frame.add_paragraph()
            p.text = clean_line
            p.level = min(level, 3)
            p.font.size = Pt(18) if level == 0 else Pt(16)
    
    return slide

def add_title_slide(prs, title, subtitle=""):
    """添加标题幻灯片"""
    slide_layout = prs.slide_layouts[0]  # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]
    
    title_shape.text = title
    subtitle_shape.text = subtitle
    
    return slide

def markdown_to_pptx(md_content, output_file):
    """将Markdown转换为PPTX"""
    # 创建演示文稿
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # 解析Markdown
    slides_data = parse_markdown(md_content)
    
    for i, slide_text in enumerate(slides_data):
        title, content = extract_title_and_content(slide_text)
        
        if not title:
            title = f"Slide {i+1}"
        
        if i == 0:
            # 第一张幻灯片作为标题幻灯片
            subtitle = '\n'.join(content[:2]) if content else ""
            add_title_slide(prs, title, subtitle)
        else:
            add_slide_with_title_and_content(prs, title, content)
    
    # 保存文件
    prs.save(output_file)
    print(f"✅ 转换完成！PPT已保存为: {output_file}")

if __name__ == "__main__":
    # 读取Markdown文件
    with open('presentation.md', 'r', encoding='utf-8') as f:
        md_content = f.read()
    
    # 转换为PPTX
    markdown_to_pptx(md_content, '长沙大学生消费行为调查.pptx')
